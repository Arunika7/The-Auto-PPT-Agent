"""
PowerPoint MCP Server — Multi-Agent V3
---------------------------------------
Design-aware MCP tool server with:
- 6 predefined theme palettes (space, business, education, tech, nature, medical)
- 4 layout engines (text_only, text_left_image_right, image_background, title_only)
- Mermaid diagram → PNG rendering via mermaid.ink API
- Internet image downloading and embedding
"""

from mcp.server.fastmcp import FastMCP
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os
import requests
import base64
from io import BytesIO

# ──────────────────────────────────────────────
# DESIGN SYSTEM: Theme Palettes
# ──────────────────────────────────────────────
THEMES = {
    "space":     {"bg": "#0b0c2a", "title": "#e879f9", "text": "#e0e7ff", "accent": "#a855f7"},
    "business":  {"bg": "#0f172a", "title": "#38bdf8", "text": "#e2e8f0", "accent": "#0ea5e9"},
    "education": {"bg": "#1c1917", "title": "#fbbf24", "text": "#fef3c7", "accent": "#f59e0b"},
    "tech":      {"bg": "#020617", "title": "#34d399", "text": "#d1fae5", "accent": "#10b981"},
    "nature":    {"bg": "#0a1628", "title": "#4ade80", "text": "#dcfce7", "accent": "#16a34a"},
    "medical":   {"bg": "#1a0f2e", "title": "#c4b5fd", "text": "#ede9fe", "accent": "#8b5cf6"},
}

def _hex_to_rgb(hex_str: str) -> RGBColor:
    """Convert hex color string to RGBColor object."""
    h = hex_str.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def _resolve_theme(theme: str) -> dict:
    """Retrieve theme colors by name, defaulting to 'space'."""
    return THEMES.get(theme.lower().strip(), THEMES["space"])

def _set_background(slide, hex_color: str):
    """Apply a solid background color to a slide."""
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = _hex_to_rgb(hex_color)

def _download_image(url: str) -> BytesIO | None:
    """Download an image from a URL and return it as a BytesIO stream."""
    try:
        # Basic validation for URL
        if not url.startswith("http"): return None
        # Fetch image with timeout and user-agent
        resp = requests.get(url, timeout=10, headers={"User-Agent": "Mozilla/5.0"})
        resp.raise_for_status()
        # Filter out very small images (often tracking pixels or icons)
        if len(resp.content) < 500: return None
        return BytesIO(resp.content)
    except Exception:
        return None

def _render_mermaid_png(mermaid_code: str) -> BytesIO | None:
    """Render Mermaid diagram code to a PNG via the mermaid.ink API."""
    try:
        # Base64 encode the Mermaid code for the API URL
        encoded = base64.urlsafe_b64encode(mermaid_code.encode("utf-8")).decode("ascii")
        url = f"https://mermaid.ink/img/{encoded}?bgColor=!0f0235"
        resp = requests.get(url, timeout=15)
        resp.raise_for_status()
        return BytesIO(resp.content)
    except Exception:
        return None

# ──────────────────────────────────────────────
# MCP SERVER INITIALIZATION
# ──────────────────────────────────────────────
mcp = FastMCP("PowerPoint MCP Server")
_current_prs = None

@mcp.tool()
def create_presentation() -> str:
    global _current_prs
    _current_prs = Presentation()
    _current_prs.slide_width = Inches(13.333)
    _current_prs.slide_height = Inches(7.5)
    return "New blank 16:9 presentation created in memory."

@mcp.tool()
def add_title_slide(title: str, subtitle: str, theme: str = "space") -> str:
    """Add a stylized title slide to the presentation."""
    global _current_prs
    if _current_prs is None: return "Error: No presentation exists."
    
    # Resolve theme colors
    colors = _resolve_theme(theme)
    # Add a blank slide layout (layout index 6 is often blank)
    slide = _current_prs.slides.add_slide(_current_prs.slide_layouts[6])
    _set_background(slide, colors["bg"])
    
    # Add and style the title text box
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.3), Inches(1.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Segoe UI"
    p.font.size = Pt(54)
    p.font.color.rgb = _hex_to_rgb(colors["title"])
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Add and style the subtitle text box
    sub_box = slide.shapes.add_textbox(Inches(2), Inches(4.2), Inches(9.3), Inches(1))
    tf2 = sub_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.name = "Segoe UI"
    p2.font.size = Pt(24)
    p2.font.color.rgb = _hex_to_rgb(colors["text"])
    p2.alignment = PP_ALIGN.CENTER
    
    # Add a thin accent bar for visual flair
    accent_bar = slide.shapes.add_shape(1, Inches(4.5), Inches(3.9), Inches(4.3), Inches(0.06))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = _hex_to_rgb(colors["accent"])
    accent_bar.line.fill.background()
    return f"Title slide '{title}' added."

@mcp.tool()
def add_slide(title: str, bullet_points: list[str], layout_type: str = "text_only",
              image_url: str = "", theme: str = "space") -> str:
    global _current_prs
    if _current_prs is None: return "Error: No presentation exists."
    
    colors = _resolve_theme(theme)
    slide = _current_prs.slides.add_slide(_current_prs.slide_layouts[6])
    _set_background(slide, colors["bg"])
    
    img_stream = _download_image(image_url) if image_url else None
    img_status = " (Visual embedded)" if img_stream else (" (Image failed)" if image_url else "")

    # Handle different layout types
    if layout_type == "title_only":
        # Centered large title for section headers
        box = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.3), Inches(2))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Segoe UI"
        p.font.size = Pt(48)
        p.font.color.rgb = _hex_to_rgb(colors["title"])
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        return f"Slide '{title}' added."

    if layout_type == "image_background" and img_stream:
        # Full-screen background image with a semi-transparent text overlay at the bottom
        slide.shapes.add_picture(img_stream, Inches(0), Inches(0), width=Inches(13.333), height=Inches(7.5))
        overlay = slide.shapes.add_shape(1, Inches(0), Inches(4.5), Inches(13.333), Inches(3))
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)
        # Apply transparency to the overlay shape
        from pptx.oxml.ns import qn
        srgb = overlay.fill._fill.find(qn("a:solidFill")).find(qn("a:srgbClr"))
        if srgb is not None: srgb.append(srgb.makeelement(qn("a:alpha"), {"val": "45000"}))
        overlay.line.fill.background()
        # Add title on top of the overlay
        t_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.7), Inches(11.7), Inches(1))
        t_box.text_frame.paragraphs[0].text = title
        t_box.text_frame.paragraphs[0].font.size = Pt(40)
        t_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        t_box.text_frame.paragraphs[0].font.bold = True
        # Add bullet points on top of the overlay
        b_box = slide.shapes.add_textbox(Inches(0.8), Inches(5.6), Inches(11.7), Inches(1.5))
        for i, bp in enumerate(bullet_points[:4]):
            p = b_box.text_frame.paragraphs[0] if i == 0 else b_box.text_frame.add_paragraph()
            p.text = f"  •  {bp}"
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(230, 230, 230)
        return f"Slide '{title}' added{img_status}."

    if layout_type == "text_left_image_right" and img_stream:
        # Split layout: Title at top, Text on left, Image on right
        t_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(1))
        t_box.text_frame.paragraphs[0].text = title
        t_box.text_frame.paragraphs[0].font.size = Pt(38)
        t_box.text_frame.paragraphs[0].font.color.rgb = _hex_to_rgb(colors["title"])
        t_box.text_frame.paragraphs[0].font.bold = True
        b_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(6.2), Inches(5))
        for i, bp in enumerate(bullet_points[:4]):
            p = b_box.text_frame.paragraphs[0] if i == 0 else b_box.text_frame.add_paragraph()
            p.text = f"•  {bp}"
            p.font.size = Pt(22)
            p.font.color.rgb = _hex_to_rgb(colors["text"])
        # Add the image to the right side of the slide
        try: slide.shapes.add_picture(img_stream, Inches(7.5), Inches(1.5), width=Inches(5.3), height=Inches(4.8))
        except Exception: pass
        return f"Slide '{title}' added{img_status}."

    # Default text_only
    t_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(1))
    t_box.text_frame.paragraphs[0].text = title
    t_box.text_frame.paragraphs[0].font.size = Pt(38)
    t_box.text_frame.paragraphs[0].font.color.rgb = _hex_to_rgb(colors["title"])
    t_box.text_frame.paragraphs[0].font.bold = True
    b_box = slide.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(10.9), Inches(4.5))
    for i, bp in enumerate(bullet_points[:4]):
        p = b_box.text_frame.paragraphs[0] if i == 0 else b_box.text_frame.add_paragraph()
        p.text = f"•  {bp}"
        p.font.size = Pt(24)
        p.font.color.rgb = _hex_to_rgb(colors["text"])
    if img_stream:
        try: slide.shapes.add_picture(img_stream, Inches(8.5), Inches(4.5), height=Inches(2.5))
        except Exception: pass
    return f"Slide '{title}' added{img_status}."

@mcp.tool()
def add_diagram_slide(title: str, mermaid_code: str, theme: str = "space") -> str:
    """Add a slide containing a Mermaid diagram rendered as an image."""
    global _current_prs
    if _current_prs is None: return "Error: No presentation exists."
    colors = _resolve_theme(theme)
    # Render the Mermaid code to an image stream
    img_stream = _render_mermaid_png(mermaid_code)
    if not img_stream: return "Error: Mermaid rendering failed."
    # Create the slide and set background
    slide = _current_prs.slides.add_slide(_current_prs.slide_layouts[6])
    _set_background(slide, colors["bg"])
    # Add title text
    t_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(11.7), Inches(1))
    t_box.text_frame.paragraphs[0].text = title
    t_box.text_frame.paragraphs[0].font.size = Pt(36)
    t_box.text_frame.paragraphs[0].font.color.rgb = _hex_to_rgb(colors["title"])
    t_box.text_frame.paragraphs[0].font.bold = True
    # Add the rendered diagram image, centered
    slide.shapes.add_picture(img_stream, Inches(1.5), Inches(1.5), width=Inches(10.3), height=Inches(5.2))
    return f"Diagram slide '{title}' added."

@mcp.tool()
def save_presentation(filename: str) -> str:
    """Save the in-memory presentation to a file."""
    global _current_prs
    if _current_prs is None: return "Error: No presentation exists."
    # Save the file with the specified name
    _current_prs.save("output_presentation.pptx")
    return "Presentation saved to output_presentation.pptx"

if __name__ == "__main__":
    mcp.run()
